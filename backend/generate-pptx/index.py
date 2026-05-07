import json
import base64
import io
import os
import uuid
import boto3
import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

NAVY = RGBColor(0x0D, 0x2D, 0x5E)
GOLD = RGBColor(0xB8, 0x96, 0x6E)
GRAY = RGBColor(0x4A, 0x55, 0x68)
GRAY_MID = RGBColor(0x71, 0x80, 0x96)
GRAY_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xEE, 0xF2, 0xF8)

IMAGES = {
    "templeClose": "https://cdn.poehali.dev/projects/1ccda9f4-d47a-42c5-9008-9a33cfe198d2/bucket/1f0e4b67-dff5-4d78-be41-f05802f1dd19.png",
    "panoramaWindows": "https://cdn.poehali.dev/projects/1ccda9f4-d47a-42c5-9008-9a33cfe198d2/bucket/3be4f93d-f0ca-41cd-b399-17e9451b9e49.png",
    "panoramaWinter": "https://cdn.poehali.dev/projects/1ccda9f4-d47a-42c5-9008-9a33cfe198d2/bucket/bf267a01-a5ac-4d36-a4a7-ab294422f50b.png",
    "bathroom": "https://cdn.poehali.dev/projects/1ccda9f4-d47a-42c5-9008-9a33cfe198d2/bucket/8622ab01-8fe2-45f9-a250-48c748a58032.png",
    "livingRoom": "https://cdn.poehali.dev/projects/1ccda9f4-d47a-42c5-9008-9a33cfe198d2/bucket/56c77e9d-d853-4cf6-afe6-7b42e04f0917.png",
    "bedroom": "https://cdn.poehali.dev/projects/1ccda9f4-d47a-42c5-9008-9a33cfe198d2/bucket/16ba467b-43f4-406c-a725-9ec119085861.png",
    "bedroom2": "https://cdn.poehali.dev/projects/1ccda9f4-d47a-42c5-9008-9a33cfe198d2/bucket/159f2f22-e971-4f8b-a42c-245fcccb9466.png",
    "triumphArch": "https://cdn.poehali.dev/projects/1ccda9f4-d47a-42c5-9008-9a33cfe198d2/bucket/fe85bea0-00d7-483f-a874-c7f4ac88ac12.png",
    "facade": "https://cdn.poehali.dev/projects/1ccda9f4-d47a-42c5-9008-9a33cfe198d2/bucket/1912d62f-a223-4e41-9ec6-02b94b5d6cab.png",
}


def fetch_image(url: str) -> io.BytesIO:
    resp = requests.get(url, timeout=15)
    resp.raise_for_status()
    return io.BytesIO(resp.content)


def add_rect(slide, x, y, w, h, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_alpha(shape, alpha_val: int):
    """Add alpha (transparency) to a solid-filled shape. alpha_val 0..100000."""
    sp = shape.fill._xPr
    solidFill = sp.find(qn('a:solidFill'))
    if solidFill is None:
        return
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        return
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', str(alpha_val))


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, font="Montserrat", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_accent_line(slide, x, y, w=Inches(0.6)):
    add_rect(slide, x, y, w, Emu(38000), GOLD)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    try:
        img = fetch_image(IMAGES["panoramaWindows"])
        slide.shapes.add_picture(img, 0, 0, width=sw, height=sh)
    except Exception:
        pass

    overlay = add_rect(slide, 0, 0, sw, sh, NAVY)
    add_alpha(overlay, 78000)

    sidebar_x = sw - Inches(2.6)
    add_rect(slide, sidebar_x, 0, Inches(2.6), sh, NAVY)

    add_text(slide, Inches(0.7), Inches(2.0), Inches(7), Inches(0.3),
             "ЭКСКЛЮЗИВНОЕ ПРЕДЛОЖЕНИЕ · 2026",
             size=11, bold=True, color=GOLD)

    add_text(slide, Inches(0.7), Inches(2.4), Inches(7.5), Inches(1.2),
             "Victory Park", size=54, color=WHITE, font="Cormorant Garamond")
    add_text(slide, Inches(0.7), Inches(3.3), Inches(7.5), Inches(1.2),
             "Residences", size=54, color=GOLD, font="Cormorant Garamond", italic=True)

    add_text(slide, Inches(0.7), Inches(4.4), Inches(7), Inches(0.5),
             "4-комнатная видовая квартира 179,08 м²", size=14, color=WHITE)
    add_text(slide, Inches(0.7), Inches(4.7), Inches(7), Inches(0.5),
             "с панорамным видом на Парк Победы", size=14, color=WHITE)
    add_text(slide, Inches(0.7), Inches(5.1), Inches(7), Inches(0.4),
             "Премиальный ремонт в стиле Неодеко · никто не проживал",
             size=11, italic=True, color=GOLD)

    stats = [("179", "м² общая"), ("60", "м² гостиная"), ("10/14", "этаж")]
    for i, (val, label) in enumerate(stats):
        x = Inches(0.7 + i * 1.8)
        add_text(slide, x, Inches(5.7), Inches(1.6), Inches(0.7),
                 val, size=32, bold=True, color=WHITE, font="Cormorant Garamond")
        add_text(slide, x, Inches(6.3), Inches(1.6), Inches(0.3),
                 label.upper(), size=8, color=WHITE)

    add_accent_line(slide, sidebar_x + Inches(0.4), Inches(1.5))
    add_text(slide, sidebar_x + Inches(0.4), Inches(1.7), Inches(2.0), Inches(0.3),
             "КЛЮЧЕВОЕ", size=9, bold=True, color=WHITE)
    add_text(slide, sidebar_x + Inches(0.4), Inches(2.1), Inches(2.0), Inches(0.4),
             "Видовой лот", size=18, bold=True, color=WHITE, font="Cormorant Garamond")
    add_text(slide, sidebar_x + Inches(0.4), Inches(2.6), Inches(2.0), Inches(1.2),
             "Один из лучших панорамных видов в Москве — на Парк Победы и центр города",
             size=9, color=WHITE)

    add_text(slide, sidebar_x + Inches(0.4), Inches(5.0), Inches(2.0), Inches(0.3),
             "СТАТУС", size=9, bold=True, color=WHITE)
    add_text(slide, sidebar_x + Inches(0.4), Inches(5.3), Inches(2.0), Inches(0.4),
             "Дом сдан · АПП подписан", size=11, bold=True, color=WHITE)

    add_rect(slide, sidebar_x + Inches(0.4), Inches(5.9), Inches(1.8), Inches(0.4), GOLD)
    add_text(slide, sidebar_x + Inches(0.4), Inches(5.97), Inches(1.8), Inches(0.3),
             "НОВАЯ КВАРТИРА", size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_about(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, sh, WHITE)

    add_accent_line(slide, Inches(0.6), Inches(0.5))
    add_text(slide, Inches(0.6), Inches(0.65), Inches(6), Inches(0.3),
             "СЛАЙД 02 · О КВАРТИРЕ", size=10, bold=True, color=GRAY_MID)

    add_rect(slide, sw - Inches(2.6), Inches(0.5), Inches(2.0), Inches(0.4), NAVY)
    add_text(slide, sw - Inches(2.6), Inches(0.57), Inches(2.0), Inches(0.3),
             "СВОБОДНАЯ ПРОДАЖА", size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.6), Inches(1.05), Inches(11), Inches(0.7),
             "Просторная резиденция в стиле Неодеко",
             size=32, color=NAVY, font="Cormorant Garamond")
    add_text(slide, Inches(0.6), Inches(1.75), Inches(11), Inches(0.3),
             "4-КОМНАТНЫЕ АПАРТАМЕНТЫ С ВИДОМ НА ПАРК ПОБЕДЫ",
             size=10, bold=True, color=GOLD)

    try:
        img = fetch_image(IMAGES["livingRoom"])
        slide.shapes.add_picture(img, Inches(0.6), Inches(2.3),
                                 width=Inches(4.8), height=Inches(4.5))
    except Exception:
        add_rect(slide, Inches(0.6), Inches(2.3), Inches(4.8), Inches(4.5), GRAY_LIGHT)

    rx = Inches(5.7)
    add_text(slide, rx, Inches(2.3), Inches(7.6), Inches(2.0),
             "Редкий видовой лот в одном из самых престижных комплексов Москвы. "
             "4-комнатная квартира 179,08 м² на 10 этаже 2 корпуса с панорамным "
             "видом на Парк Победы и центр Москвы. Дом сдан, акт приёма-передачи "
             "подписан. Премиальный дизайнерский ремонт в стиле Неодеко — в "
             "квартире никто не проживал.",
             size=11, color=GRAY)

    specs = [("Комнат", "4"), ("Площадь", "179,08 м²"), ("Жилая", "90 м²"),
             ("Кухня", "60 м²"), ("Этаж", "10 / 14"), ("Корпус", "2")]
    spec_y = Inches(4.4)
    for i, (label, val) in enumerate(specs):
        col = i % 3
        row = i // 3
        x = rx + Inches(col * 2.55)
        y = spec_y + Inches(row * 0.8)
        add_rect(slide, x, y, Inches(2.4), Inches(0.7), BG_LIGHT)
        add_rect(slide, x, y, Emu(38000), Inches(0.7), NAVY)
        add_text(slide, x + Inches(0.15), y + Inches(0.08), Inches(2.2), Inches(0.25),
                 label.upper(), size=8, bold=True, color=GRAY_MID)
        add_text(slide, x + Inches(0.15), y + Inches(0.32), Inches(2.2), Inches(0.35),
                 val, size=12, bold=True, color=NAVY)

    features = [("Кухня-гостиная 60 м²", "Просторная зона для жизни"),
                ("3 мастер-спальни", "С санузлами и гардеробными"),
                ("Стиль Неодеко", "Премиальный ремонт"),
                ("Панорамные окна", "Максимум света")]
    fy = Inches(6.1)
    for i, (title, sub) in enumerate(features):
        col = i % 2
        row = i // 2
        x = rx + Inches(col * 3.85)
        y = fy + Inches(row * 0.4)
        add_rect(slide, x, y, Inches(0.3), Inches(0.3), NAVY)
        add_text(slide, x + Inches(0.4), y, Inches(3.4), Inches(0.2),
                 title, size=10, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.4), y + Inches(0.18), Inches(3.4), Inches(0.2),
                 sub, size=8, color=GRAY_MID)


def slide_location(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, sh, BG_LIGHT)

    left_w = sw - Inches(3.5)
    try:
        img = fetch_image(IMAGES["templeClose"])
        slide.shapes.add_picture(img, 0, 0, width=left_w, height=sh)
    except Exception:
        pass
    overlay = add_rect(slide, 0, 0, left_w, sh, NAVY)
    add_alpha(overlay, 70000)

    add_accent_line(slide, Inches(0.6), Inches(0.7))
    add_text(slide, Inches(0.6), Inches(0.85), Inches(7), Inches(0.3),
             "СЛАЙД 03 · ЛОКАЦИЯ", size=10, bold=True, color=WHITE)

    add_text(slide, Inches(0.6), Inches(1.3), Inches(8), Inches(0.9),
             "Лучший вид", size=42, color=WHITE, font="Cormorant Garamond")
    add_text(slide, Inches(0.6), Inches(2.1), Inches(8), Inches(0.9),
             "в столице", size=42, color=GOLD, font="Cormorant Garamond", italic=True)

    add_text(slide, Inches(0.6), Inches(3.2), Inches(7), Inches(0.3),
             "ИЗ ОКОН ВИДНО", size=10, bold=True, color=GOLD)

    views = ["Музей и Монумент Победы", "Храм Георгия Победоносца",
             "Фонтан «Годы войны»", "Триумфальные ворота", "Москва-Сити"]
    for i, v in enumerate(views):
        y = Inches(3.55 + i * 0.32)
        add_text(slide, Inches(0.6), y, Inches(0.2), Inches(0.25),
                 "•", size=14, bold=True, color=GOLD)
        add_text(slide, Inches(0.85), y, Inches(7), Inches(0.25),
                 v, size=12, color=WHITE)

    locs = [("Напротив", "ПАРК ПОБЕДЫ"), ("15 мин", "ДО КРЕМЛЯ"),
            ("10 мин", "ДО МОСКВА-СИТИ"), ("5 мин", "ДО ЛУЖНИКОВ")]
    ly = Inches(5.45)
    for i, (val, label) in enumerate(locs):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 3.0)
        y = ly + Inches(row * 0.85)
        add_rect(slide, x, y, Inches(2.85), Inches(0.75), GRAY)
        add_rect(slide, x, y, Emu(38000), Inches(0.75), GOLD)
        add_text(slide, x + Inches(0.2), y + Inches(0.08), Inches(2.6), Inches(0.4),
                 val, size=18, bold=True, color=WHITE, font="Cormorant Garamond")
        add_text(slide, x + Inches(0.2), y + Inches(0.45), Inches(2.6), Inches(0.25),
                 label, size=8, bold=True, color=WHITE)

    sx = sw - Inches(3.5)
    add_rect(slide, sx, 0, Inches(3.5), sh, WHITE)
    add_accent_line(slide, sx + Inches(0.4), Inches(0.7))
    add_text(slide, sx + Inches(0.4), Inches(0.95), Inches(3), Inches(0.7),
             "Инфраструктура ЖК", size=22, color=NAVY, font="Cormorant Garamond")
    add_text(slide, sx + Inches(0.4), Inches(1.7), Inches(3), Inches(0.4),
             "Закрытая территория с дизайнерским ландшафтом",
             size=9, color=GRAY_MID)

    amenities = ["Фитнес-центр с бассейном", "Wellness и салон красоты",
                 "Аэройога и пилатес", "Закрытая озеленённая территория",
                 "Детский сад и площадки", "Подземный паркинг и автомойка"]
    for i, a in enumerate(amenities):
        y = Inches(2.4 + i * 0.65)
        add_rect(slide, sx + Inches(0.4), y, Inches(3), Inches(0.55), BG_LIGHT)
        add_rect(slide, sx + Inches(0.4), y, Inches(0.5), Inches(0.55), NAVY)
        add_text(slide, sx + Inches(1.0), y + Inches(0.15), Inches(2.4), Inches(0.3),
                 a, size=10, bold=True, color=NAVY)


def slide_gallery(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, sh, BG_LIGHT)

    add_accent_line(slide, Inches(0.6), Inches(0.5))
    add_text(slide, Inches(0.6), Inches(0.7), Inches(8), Inches(0.7),
             "Галерея фотографий", size=32, color=NAVY, font="Cormorant Garamond")
    add_text(slide, sw - Inches(3.5), Inches(0.85), Inches(3), Inches(0.3),
             "СЛАЙД 04 · СЪЁМКА 24.04.2026", size=10, bold=True, color=GRAY_MID,
             align=PP_ALIGN.RIGHT)

    photos = [
        (IMAGES["livingRoom"], "Гостиная · панорамное остекление"),
        (IMAGES["panoramaWindows"], "Изогнутые панорамные окна"),
        (IMAGES["panoramaWinter"], "Вид зимой · Храм и центр Москвы"),
        (IMAGES["bathroom"], "Мастер-санузел · мрамор"),
        (IMAGES["bedroom"], "Спальня с панорамным окном"),
        (IMAGES["bedroom2"], "Спальня · натуральные материалы"),
        (IMAGES["templeClose"], "Вид · Храм Георгия Победоносца"),
        (IMAGES["triumphArch"], "Триумфальные ворота — рядом"),
        (IMAGES["facade"], "Victory Park Residences · фасад"),
    ]

    cols, rows = 3, 3
    grid_x = Inches(0.6)
    grid_y = Inches(1.8)
    cell_w = Inches(4.0)
    cell_h = Inches(1.7)
    gap = Inches(0.15)
    for i, (url, caption) in enumerate(photos):
        col = i % cols
        row = i // cols
        x = grid_x + col * (cell_w + gap)
        y = grid_y + row * (cell_h + gap)
        try:
            img = fetch_image(url)
            slide.shapes.add_picture(img, x, y, width=cell_w, height=cell_h)
        except Exception:
            add_rect(slide, x, y, cell_w, cell_h, GRAY_LIGHT)
        cap_h = Inches(0.3)
        add_rect(slide, x, y + cell_h - cap_h, cell_w, cap_h, NAVY)
        add_text(slide, x + Inches(0.1), y + cell_h - cap_h + Inches(0.05),
                 cell_w - Inches(0.2), cap_h, caption, size=8, bold=True, color=WHITE)


def slide_parking(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, sh, WHITE)

    right_w = Inches(4.5)
    rx = sw - right_w
    try:
        img = fetch_image(IMAGES["facade"])
        slide.shapes.add_picture(img, rx, 0, width=right_w, height=sh)
    except Exception:
        add_rect(slide, rx, 0, right_w, sh, GRAY_LIGHT)

    add_accent_line(slide, Inches(0.6), Inches(0.7))
    add_text(slide, Inches(0.6), Inches(0.9), Inches(7), Inches(0.3),
             "СЛАЙД 05 · ДОПОЛНИТЕЛЬНО", size=10, bold=True, color=GRAY_MID)

    add_text(slide, Inches(0.6), Inches(1.4), Inches(8), Inches(0.9),
             "3 машино-места", size=38, color=NAVY, font="Cormorant Garamond")
    add_text(slide, Inches(0.6), Inches(2.15), Inches(8), Inches(0.9),
             "одним блоком", size=38, color=GOLD, font="Cormorant Garamond", italic=True)

    add_text(slide, Inches(0.6), Inches(3.1), Inches(7.5), Inches(1.2),
             "В продаже также три машино-места, расположенные рядом друг с другом "
             "одним блоком — для вашего парка автомобилей. Подземный охраняемый "
             "паркинг с прямым доступом с любого этажа комплекса.",
             size=11, color=GRAY)

    add_rect(slide, Inches(0.6), Inches(4.5), Inches(3.6), Inches(1.4), NAVY)
    add_text(slide, Inches(0.8), Inches(4.6), Inches(3.4), Inches(0.3),
             "3 МАШИНО-МЕСТА", size=9, bold=True, color=WHITE)
    add_text(slide, Inches(0.8), Inches(4.9), Inches(3.4), Inches(0.6),
             "50 млн ₽", size=28, bold=True, color=WHITE, font="Cormorant Garamond")
    add_text(slide, Inches(0.8), Inches(5.5), Inches(3.4), Inches(0.3),
             "Не входит в стоимость квартиры", size=8, color=WHITE)

    add_rect(slide, Inches(4.4), Inches(4.5), Inches(3.6), Inches(1.4), BG_LIGHT, GOLD)
    add_text(slide, Inches(4.6), Inches(4.6), Inches(3.4), Inches(0.3),
             "ОПЦИОНАЛЬНО", size=9, bold=True, color=GOLD)
    add_text(slide, Inches(4.6), Inches(4.9), Inches(3.4), Inches(0.9),
             "Можно купить только квартиру — на машино-места уже есть покупатель",
             size=11, bold=True, color=NAVY)

    bullets = ["Охраняемый подземный паркинг 24/7",
               "Автомойка на территории комплекса",
               "Прямой доступ с любого этажа"]
    for i, b in enumerate(bullets):
        y = Inches(6.1 + i * 0.35)
        add_rect(slide, Inches(0.6), y, Inches(0.25), Inches(0.25), NAVY)
        add_text(slide, Inches(1.0), y, Inches(7), Inches(0.3),
                 b, size=10, bold=True, color=NAVY)


def slide_contacts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, sh, NAVY)

    add_accent_line(slide, Inches(0.6), Inches(0.7))
    add_text(slide, Inches(0.6), Inches(0.9), Inches(7), Inches(0.3),
             "СЛАЙД 06 · СВЯЗАТЬСЯ", size=10, bold=True, color=WHITE)

    add_text(slide, Inches(0.6), Inches(1.3), Inches(8), Inches(0.9),
             "Готовы показать", size=38, color=WHITE, font="Cormorant Garamond")
    add_text(slide, Inches(0.6), Inches(2.05), Inches(8), Inches(0.9),
             "квартиру лично", size=38, color=GOLD, font="Cormorant Garamond", italic=True)

    add_text(slide, Inches(0.6), Inches(3.1), Inches(8), Inches(1.0),
             "Звоните или пишите — отвечу максимально быстро. Расскажу подробности, "
             "согласую удобное время просмотра и пришлю полный пакет документов.",
             size=11, color=WHITE)

    contacts = [("ТЕЛЕФОН", "+7 (___) ___-__-__"),
                ("TELEGRAM / WHATSAPP", "@username"),
                ("EMAIL", "info@example.ru")]
    for i, (label, val) in enumerate(contacts):
        y = Inches(4.4 + i * 0.7)
        add_rect(slide, Inches(0.6), y, Inches(0.5), Inches(0.5), GRAY)
        add_text(slide, Inches(1.25), y, Inches(7), Inches(0.25),
                 label, size=8, bold=True, color=WHITE)
        add_text(slide, Inches(1.25), y + Inches(0.22), Inches(7), Inches(0.4),
                 val, size=14, bold=True, color=WHITE)

    sx = sw - Inches(4.0)
    panel = add_rect(slide, sx, 0, Inches(4.0), sh, GRAY)
    add_alpha(panel, 50000)

    add_text(slide, sx + Inches(0.4), Inches(0.7), Inches(3.5), Inches(0.3),
             "РЕЗЮМЕ ОБЪЕКТА", size=10, bold=True, color=GOLD)

    summary = [("Объект", "4-комн. квартира"), ("Площадь", "179,08 м²"),
               ("Этаж", "10 / 14, корпус 2"), ("Состояние", "Новая · Неодеко"),
               ("Статус", "Свободная продажа"), ("Готовность", "Ключи на руках")]
    for i, (label, val) in enumerate(summary):
        y = Inches(1.3 + i * 0.55)
        add_text(slide, sx + Inches(0.4), y, Inches(1.5), Inches(0.3),
                 label, size=9, color=WHITE)
        add_text(slide, sx + Inches(1.9), y, Inches(2.0), Inches(0.3),
                 val, size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    add_rect(slide, sx + Inches(0.4), Inches(6.0), Inches(3.2), Inches(1.0), GOLD)
    add_text(slide, sx + Inches(0.6), Inches(6.15), Inches(2.8), Inches(0.3),
             "РЕДКОЕ ПРЕДЛОЖЕНИЕ", size=9, bold=True, color=WHITE)
    add_text(slide, sx + Inches(0.6), Inches(6.4), Inches(2.8), Inches(0.5),
             "Видовых лотов такого формата на рынке единицы.",
             size=9, color=WHITE)


def build_pptx() -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_about(prs)
    slide_location(prs)
    slide_gallery(prs)
    slide_parking(prs)
    slide_contacts(prs)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def handler(event: dict, context) -> dict:
    """Генерирует .pptx презентацию Victory Park Residences и возвращает её base64 для скачивания."""
    method = event.get('httpMethod', 'GET')

    cors = {
        'Access-Control-Allow-Origin': '*',
        'Access-Control-Allow-Methods': 'GET, POST, OPTIONS',
        'Access-Control-Allow-Headers': 'Content-Type, X-User-Id, X-Auth-Token, X-Session-Id',
        'Access-Control-Max-Age': '86400',
    }

    if method == 'OPTIONS':
        return {'statusCode': 200, 'headers': cors, 'body': ''}

    try:
        data = build_pptx()

        s3 = boto3.client(
            's3',
            endpoint_url='https://bucket.poehali.dev',
            aws_access_key_id=os.environ['AWS_ACCESS_KEY_ID'],
            aws_secret_access_key=os.environ['AWS_SECRET_ACCESS_KEY'],
        )
        key = f"presentations/VictoryPark_{uuid.uuid4().hex[:8]}.pptx"
        s3.put_object(
            Bucket='files',
            Key=key,
            Body=data,
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            ContentDisposition='attachment; filename="VictoryPark_Presentation.pptx"',
        )
        url = f"https://cdn.poehali.dev/projects/{os.environ['AWS_ACCESS_KEY_ID']}/bucket/{key}"

        return {
            'statusCode': 200,
            'headers': {**cors, 'Content-Type': 'application/json'},
            'body': json.dumps({'url': url, 'filename': 'VictoryPark_Presentation.pptx'}),
        }
    except Exception as e:
        import traceback
        return {
            'statusCode': 500,
            'headers': {**cors, 'Content-Type': 'application/json'},
            'body': json.dumps({'error': str(e), 'trace': traceback.format_exc()}),
        }